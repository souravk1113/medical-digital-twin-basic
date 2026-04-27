"""Build slides/medical_digital_twin_GA.pptx.

A presentable deck that walks from the modern (foundation-model / LLM-augmented)
digital-twin landscape to a concrete proof-of-concept proposal for a digital
twin of Geographic Atrophy (GA) in ophthalmology.

Run from repo root with the project venv active:
    python slides/build_deck.py
"""
from __future__ import annotations
import datetime as dt
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------
NAVY   = RGBColor(0x0B, 0x25, 0x45)
TEAL   = RGBColor(0x1B, 0x6C, 0xA8)
ACCENT = RGBColor(0x2E, 0x8B, 0xC4)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GREY   = RGBColor(0x55, 0x5B, 0x66)
LIGHT  = RGBColor(0xE9, 0xF0, 0xF6)
PALE   = RGBColor(0xF5, 0xF8, 0xFB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RED    = RGBColor(0xC0, 0x39, 0x2B)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "slides" / "medical_digital_twin_GA.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_text(shape, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def _bullets(shape, items, size=16, color=INK, indent_px=0):
    """items is a list of (text, level) tuples, or strings (level=0)."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    first = True
    for entry in items:
        if isinstance(entry, tuple):
            text, level = entry
        else:
            text, level = entry, 0
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        bullet = "•  " if level == 0 else "–  "
        run = p.add_run()
        run.text = bullet + text
        run.font.name = "Calibri"
        run.font.size = Pt(size if level == 0 else max(12, size - 2))
        run.font.color.rgb = color if level == 0 else GREY
        run.font.bold = False
        p.space_after = Pt(4)


def _add_rect(slide, x, y, w, h, fill=LIGHT, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _set_fill(shp, fill)
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    return shp


def _add_arrow(slide, x1, y1, x2, y2, color=TEAL, weight=2.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.W = self.prs.slide_width
        self.H = self.prs.slide_height
        self.section_idx = 0
        self.total_sections = 7

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    # ---- specific slide types ----
    def title_slide(self, title, subtitle, footer):
        s = self._blank()
        # left navy band
        _add_rect(s, 0, 0, Inches(0.6), self.H, fill=NAVY)
        # accent strip
        _add_rect(s, Inches(0.9), Inches(2.6), Inches(0.7), Inches(0.06), fill=TEAL)
        # title
        tb = s.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(11.5), Inches(1.6))
        _add_text(tb, title, size=40, bold=True, color=NAVY)
        # subtitle
        sb = s.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.5), Inches(0.9))
        _add_text(sb, subtitle, size=20, color=GREY)
        # footer
        fb = s.shapes.add_textbox(Inches(1.0), Inches(6.6), Inches(11.5), Inches(0.4))
        _add_text(fb, footer, size=12, color=GREY)
        # decorative dots
        for i, c in enumerate([TEAL, ACCENT, NAVY]):
            d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(11.4 + i * 0.45), Inches(6.55),
                                   Inches(0.18), Inches(0.18))
            _set_fill(d, c)

    def section_divider(self, num, title, subtitle=""):
        s = self._blank()
        # full navy background
        _add_rect(s, 0, 0, self.W, self.H, fill=NAVY)
        # giant section number
        nb = s.shapes.add_textbox(Inches(0.7), Inches(0.9), Inches(4.0), Inches(2.0))
        _add_text(nb, f"{num:02d}", size=110, bold=True, color=TEAL)
        # title
        tb = s.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(12.0), Inches(1.5))
        _add_text(tb, title, size=44, bold=True, color=WHITE)
        # subtitle
        if subtitle:
            sb = s.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(12.0), Inches(1.2))
            _add_text(sb, subtitle, size=18, color=LIGHT)
        # accent strip
        _add_rect(s, Inches(0.7), Inches(6.6), Inches(1.5), Inches(0.06), fill=TEAL)

    def _content_header(self, slide, title, eyebrow=None):
        # header bar
        _add_rect(slide, 0, 0, self.W, Inches(1.05), fill=NAVY)
        _add_rect(slide, 0, Inches(1.05), self.W, Inches(0.05), fill=TEAL)
        if eyebrow:
            eb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18),
                                          Inches(12.0), Inches(0.3))
            _add_text(eb, eyebrow.upper(), size=10, bold=True, color=ACCENT)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.42),
                                      Inches(12.0), Inches(0.6))
        _add_text(tb, title, size=24, bold=True, color=WHITE)

    def content(self, title, bullets, eyebrow=None, footnote=None):
        s = self._blank()
        self._content_header(s, title, eyebrow)
        # body
        body = s.shapes.add_textbox(Inches(0.5), Inches(1.4),
                                    Inches(12.3), Inches(5.6))
        _bullets(body, bullets)
        if footnote:
            fb = s.shapes.add_textbox(Inches(0.5), Inches(7.05),
                                      Inches(12.3), Inches(0.35))
            _add_text(fb, footnote, size=10, color=GREY)
        return s

    def two_column(self, title, left_head, left_bullets, right_head,
                   right_bullets, eyebrow=None, footnote=None):
        s = self._blank()
        self._content_header(s, title, eyebrow)
        # left column header
        lh = s.shapes.add_textbox(Inches(0.5), Inches(1.35),
                                  Inches(6.0), Inches(0.4))
        _add_text(lh, left_head, size=15, bold=True, color=NAVY)
        # left body
        lb = s.shapes.add_textbox(Inches(0.5), Inches(1.8),
                                  Inches(6.1), Inches(5.2))
        _bullets(lb, left_bullets, size=14)
        # divider
        _add_rect(s, Inches(6.65), Inches(1.4), Inches(0.02), Inches(5.6),
                  fill=LIGHT)
        # right column header
        rh = s.shapes.add_textbox(Inches(6.85), Inches(1.35),
                                  Inches(6.0), Inches(0.4))
        _add_text(rh, right_head, size=15, bold=True, color=NAVY)
        # right body
        rb = s.shapes.add_textbox(Inches(6.85), Inches(1.8),
                                  Inches(6.1), Inches(5.2))
        _bullets(rb, right_bullets, size=14)
        if footnote:
            fb = s.shapes.add_textbox(Inches(0.5), Inches(7.05),
                                      Inches(12.3), Inches(0.35))
            _add_text(fb, footnote, size=10, color=GREY)
        return s

    def callout(self, title, bullets, callout_text, eyebrow=None):
        s = self._blank()
        self._content_header(s, title, eyebrow)
        # bullets on left
        body = s.shapes.add_textbox(Inches(0.5), Inches(1.4),
                                    Inches(7.5), Inches(5.6))
        _bullets(body, bullets, size=15)
        # callout box on right
        box = _add_rect(s, Inches(8.4), Inches(1.55),
                        Inches(4.5), Inches(4.5), fill=PALE, line=ACCENT)
        eb = s.shapes.add_textbox(Inches(8.6), Inches(1.7),
                                  Inches(4.2), Inches(0.4))
        _add_text(eb, "KEY TAKEAWAY", size=10, bold=True, color=ACCENT)
        cb = s.shapes.add_textbox(Inches(8.6), Inches(2.1),
                                  Inches(4.2), Inches(3.8))
        _add_text(cb, callout_text, size=14, color=NAVY)
        return s

    def table_slide(self, title, headers, rows, eyebrow=None,
                    col_widths=None, footnote=None):
        s = self._blank()
        self._content_header(s, title, eyebrow)
        n_cols = len(headers)
        n_rows = len(rows) + 1
        left = Inches(0.5)
        top = Inches(1.45)
        width = Inches(12.3)
        height = Inches(0.5 * n_rows + 0.3)
        tbl_shape = s.shapes.add_table(n_rows, n_cols, left, top, width, height)
        tbl = tbl_shape.table
        if col_widths:
            for i, w in enumerate(col_widths):
                tbl.columns[i].width = Inches(w)
        # header
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            cell.text_frame.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = h
            r.font.bold = True
            r.font.size = Pt(13)
            r.font.color.rgb = WHITE
            r.font.name = "Calibri"
            cell.text_frame.margin_top = Inches(0.05)
            cell.text_frame.margin_bottom = Inches(0.05)
        # rows
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else PALE
                cell.text_frame.text = ""
                p = cell.text_frame.paragraphs[0]
                r = p.add_run()
                r.text = str(val)
                r.font.size = Pt(11)
                r.font.color.rgb = INK
                r.font.name = "Calibri"
                cell.text_frame.margin_top = Inches(0.04)
                cell.text_frame.margin_bottom = Inches(0.04)
        if footnote:
            fb = s.shapes.add_textbox(Inches(0.5), Inches(7.05),
                                      Inches(12.3), Inches(0.35))
            _add_text(fb, footnote, size=10, color=GREY)
        return s

    def architecture_slide(self, title, eyebrow=None):
        """A bespoke 5-layer architecture diagram."""
        s = self._blank()
        self._content_header(s, title, eyebrow)
        # patient box (left)
        x_p, y_p, w_p, h_p = Inches(0.6), Inches(2.5), Inches(2.2), Inches(2.5)
        _add_rect(s, x_p, y_p, w_p, h_p, fill=LIGHT, line=NAVY)
        tp = s.shapes.add_textbox(x_p, y_p + Inches(0.1), w_p, Inches(0.6))
        _add_text(tp, "PHYSICAL\nPATIENT", size=14, bold=True,
                  color=NAVY, align=PP_ALIGN.CENTER)
        info = s.shapes.add_textbox(x_p, y_p + Inches(0.9), w_p, Inches(1.6))
        _bullets(info, ["FAF, OCT", "BCVA, LL-VA", "Genetics", "Demographics"],
                 size=11)
        # 5-layer stack (right)
        layers = [
            ("1  DATA",        "FAF · OCT · EHR · genetics → FHIR / OMOP", LIGHT),
            ("2  STATE",       "Lesion area, focality, FAF pattern, BCVA, time", LIGHT),
            ("3  MODEL",       "sqrt growth law + ML residual + LLM orchestrator", LIGHT),
            ("4  SIMULATION",  "What-if drugs · counterfactuals · in-silico trials", LIGHT),
            ("5  DECISION/UI", "Clinician dashboard · trial-design panel · alerts", LIGHT),
        ]
        x0 = Inches(3.4); y0 = Inches(1.4); w = Inches(9.4); h = Inches(1.05)
        for i, (head, body, fill) in enumerate(layers):
            y = y0 + i * (h + Inches(0.05))
            _add_rect(s, x0, y, w, h, fill=fill, line=ACCENT)
            hb = s.shapes.add_textbox(x0 + Inches(0.15), y + Inches(0.1),
                                      Inches(2.1), Inches(0.85))
            _add_text(hb, head, size=15, bold=True, color=NAVY,
                      anchor=MSO_ANCHOR.MIDDLE)
            bb = s.shapes.add_textbox(x0 + Inches(2.3), y + Inches(0.1),
                                      Inches(7.0), Inches(0.85))
            _add_text(bb, body, size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        # arrow patient -> stack
        _add_arrow(s, x_p + w_p, Inches(3.0),
                   x0, Inches(3.0), color=TEAL, weight=2.5)
        # feedback arrow back
        _add_arrow(s, x0, Inches(6.2),
                   x_p + w_p, Inches(6.2), color=ACCENT, weight=2.0)
        fb = s.shapes.add_textbox(Inches(0.6), Inches(6.4),
                                  Inches(12.3), Inches(0.4))
        _add_text(fb, "data flows up · decisions flow down — the closed loop is what makes it a TWIN, not a model.",
                  size=11, color=GREY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# CONTENT
# ---------------------------------------------------------------------------
def build():
    d = Deck()

    # ===== Title =====
    d.title_slide(
        "A Digital Twin for Geographic Atrophy",
        "From classical medical twins to foundation-model + LLM-augmented twins, "
        "and a concrete proof-of-concept proposal in ophthalmology.",
        f"Sourav K  ·  {dt.date.today().isoformat()}  ·  internal review draft",
    )

    # ===== TL;DR =====
    d.callout(
        "Executive summary",
        [
            "Medical digital twins have moved beyond bespoke ODE models into a foundation-model + LLM-orchestrated era.",
            "Geographic Atrophy (GA), an advanced dry-AMD disease, is an unusually good fit for a twin: quantitative imaging endpoint, slow progression, partly known mechanism, two new drugs with modest effect sizes.",
            "We propose a tractable v0 PoC: a sqrt-area growth-law twin with ML personalization, a Streamlit dashboard, and an in-silico trial module — buildable on simulated cohorts calibrated to published distributions, no credentialed data needed.",
            "Roadmap escalates to AREDS2 (dbGaP), RETFound embeddings, and an LLM orchestrator.",
        ],
        "GA's quantitative FAF endpoint, slow trial timelines, and two FDA-approved drugs with partial effect make it the ophthalmology twin worth building first.",
        eyebrow="TL;DR",
    )

    # ===== Agenda =====
    d.content(
        "Agenda",
        [
            ("1.  Digital twins — recap", 0),
            ("2.  Modern paradigms: foundation models, neural ODEs, generative twins, LLM orchestration, in-silico trials", 0),
            ("3.  Geographic Atrophy primer", 0),
            ("4.  Why a GA twin?", 0),
            ("5.  Prior art in ophthalmology AI", 0),
            ("6.  Proposed architecture", 0),
            ("7.  Proof-of-concept: scope, data, methodology, validation", 0),
            ("8.  Roadmap, risks, references", 0),
        ],
        eyebrow="Roadmap",
    )

    # ====== SECTION 1 ======
    d.section_divider(1, "Digital twins — recap",
                      "Where the field came from and what the three layers are.")

    d.content(
        "What is a digital twin?",
        [
            "A virtual representation of a real-world entity that is continuously updated by data from its physical counterpart and supports simulation, prediction and decision-making in a closed feedback loop.",
            ("Origins: Grieves (2002, PLM) → NASA (2010, aerospace) → medicine from ~2018 onward.", 1),
            "Three components every twin has:",
            ("Object — patient / organ / cell / population — represented by mathematical and ML models.", 1),
            ("State — vitals, labs, images, history — a time-evolving state vector x(t).", 1),
            ("Behavior — physiology and disease progression — simulated by ODE / agent / ML / hybrid models.", 1),
            "Distinction from a regular predictive model: a twin is persistent, individualized, and updatable — a risk score collapses a patient to one number; a twin is a queryable, simulatable representation over time.",
        ],
        eyebrow="01 · Recap",
    )

    d.two_column(
        "Classical medical twins (pre-2022)",
        "Examples in production",
        [
            "Cardiac: Dassault Living Heart, Siemens CRT response twin, HeartFlow CFD on coronary CT.",
            "Endocrinology: UVa/Padova T1D simulator (FDA-recognized in-silico testing for closed-loop pumps); Twin Health metabolic twin for T2D reversal (RCT).",
            "Oncology: image-guided reaction-diffusion tumor twins (Yankeelov; Swanson).",
            "ICU: sepsis policy twins on MIMIC; ARDS ventilator twins.",
        ],
        "What they have in common",
        [
            "Mechanistic core (ODE / PDE / FEM) encoding known physiology.",
            "Per-patient calibration of a small number of parameters.",
            "Heavy domain expertise per project — not portable.",
            "Limited multimodal input — usually one or two data streams.",
            "Validation via held-out data and prospective clinical studies.",
        ],
        eyebrow="01 · Recap",
    )

    # ===== SECTION 2 — MODERN PARADIGMS =====
    d.section_divider(2, "Modern paradigms (2022 → 2026)",
                      "Foundation models, neural ODEs, generative twins, "
                      "LLM-orchestrated stacks, in-silico trials.")

    d.callout(
        "Foundation models as the perception layer",
        [
            "Self-supervised pre-training on massive medical-data corpora collapses the 'feature engineering' step.",
            "RETFound (Zhou et al., Nature 2023) — masked-autoencoder ViT, ~1.6 M retinal images (CFP + OCT), open weights. State-of-the-art on AMD progression. [1]",
            "Med-PaLM 2 (Singhal et al. 2023) — clinical-question LLM. [2]",
            "Generalist Medical AI (Moor et al., Nature 2023) — multimodal foundation model for clinical reasoning. [3]",
            "MedGemma (Google 2024) — open multimodal medical foundation model. [4]",
            "AMIE (Tu et al. 2024) — Google diagnostic-dialogue agent. [5]",
        ],
        "Twin builders increasingly inherit the perception layer for free; the value-add moves to mechanism, simulation, and the closed loop.",
        eyebrow="02 · Modern · perception",
    )

    d.content(
        "Neural ODEs and physics-informed models",
        [
            "Hybrid mechanistic + ML models close the data-efficiency gap of pure deep nets.",
            "Neural ODEs (Chen et al., NeurIPS 2018) — continuous-time hidden state, ideal for irregularly-sampled clinical time series. [6]",
            "Physics-informed neural networks — penalize violations of known PDEs in the loss; useful when mechanism is partly understood. [7]",
            "Universal differential equations — keep the mechanistic skeleton, learn the residual with an NN. [8]",
            "For GA this maps cleanly to: mechanism = sqrt-area growth law (Yehoshua, Feuer); residual = patient-specific deviations learned from imaging.",
            "Practical effect: 5–50× less data than a pure deep net for the same forecast quality, plus extrapolation guarantees from the mechanistic part.",
        ],
        eyebrow="02 · Modern · hybrid models",
    )

    d.content(
        "Generative twins and synthetic patients",
        [
            "Diffusion models for medical imaging (Kazerouni et al. 2023) — counterfactual images: 'what would this fundus look like in 24 months under treatment?' [9]",
            "Tabular synthesis (CTGAN / TVAE / tabular diffusion) — privacy-preserving cohorts; basis for synthetic control arms. [10]",
            "Generative EHR time-series (Theodorou et al., npj Digital Medicine 2023) — long synthetic patient histories that respect dependencies. [11]",
            "Important caveat: generative output can leak training data and under-represent rare phenotypes; mitigation = privacy testing, rare-class oversampling, expert review.",
            "For a GA twin: useful for visualizing projected lesions to clinicians and patients; decisions still drive off the mechanistic + ML layer.",
        ],
        eyebrow="02 · Modern · generative",
    )

    d.content(
        "LLM-orchestrated twins",
        [
            "Pattern emerging in 2024–2026 deployments: the LLM does NOT predict; it routes the clinician's natural-language question to a validated mechanistic / ML twin and explains the result.",
            ("Architecture: clinician → LLM orchestrator (tool-using, retrieval-augmented) → calls into mechanistic twin and a knowledge base (FHIR / OMOP / guideline corpus) → answers in natural language with citations.", 0),
            ("Why this shape: keeps the clinical decision auditable (mechanistic core is the decision-maker), while letting the LLM handle the open-vocabulary clinician query.", 1),
            ("Aligns with FDA AI/ML SaMD action plan and EMA reflection paper (2024) — separation of reasoning surface from decision surface. [12]", 1),
            "Downside: prompt-injection, hallucinated tool calls, and citation drift remain real; need eval harnesses and red-teaming before deployment.",
        ],
        eyebrow="02 · Modern · orchestration",
    )

    d.callout(
        "In-silico trials & synthetic control arms",
        [
            "FDA's MIDD program accepts mechanistic models in submissions for dose justification and label extension. [13]",
            "External / synthetic control arms accepted in oncology (FlatironHealth) and devices; increasingly proposed for rare-disease and slow-progressing indications. [14]",
            "Twin populations (Unlearn.AI) — generate digital twins of every control patient to reduce trial sample size by ~30%. EMA-reviewed for use in some Phase 2/3 trials, including Alzheimer's. [15]",
            "GA trials run 12–24 months and recruit slowly → enormous incentive to augment control arms.",
        ],
        "GA is one of the most attractive indications for twin-augmented control arms — slow, quantitative, with a structural endpoint.",
        eyebrow="02 · Modern · regulatory",
    )

    # ====== SECTION 3 — GA PRIMER ======
    d.section_divider(3, "Geographic Atrophy",
                      "What it is, how we measure it, what we treat it with, what's missing.")

    d.content(
        "What is Geographic Atrophy?",
        [
            "The advanced, late-stage form of dry age-related macular degeneration (AMD).",
            "Progressive, irreversible degeneration of the retinal pigment epithelium (RPE), choriocapillaris and overlying photoreceptors.",
            "Affects ~5 million people worldwide; a leading cause of irreversible central-vision loss in adults > 50. [16]",
            "Bilateral in most patients eventually; vision loss usually starts paracentrally and encroaches on the fovea.",
            "Pathophysiology: complement-mediated chronic inflammation + oxidative stress + lipofuscin accumulation → RPE death → photoreceptor loss. Genetics matters: CFH, ARMS2/HTRA1, C3.",
        ],
        eyebrow="03 · GA · primer",
    )

    d.two_column(
        "How we measure GA",
        "Imaging biomarkers",
        [
            "Fundus autofluorescence (FAF) — gold standard; lesions hypo-autofluorescent.",
            "Spectral-domain OCT — RPE atrophy, photoreceptor (EZ) loss, drusen volume, hyperreflective foci, reticular pseudodrusen.",
            "Color fundus photography — historical baseline.",
            "Microperimetry — point-wise retinal sensitivity.",
        ],
        "Quantitative endpoints",
        [
            "Lesion area in mm² — primary trial endpoint.",
            "Square-root transform (sqrt mm²/yr) linearizes growth across baseline sizes (Feuer 2013). [18]",
            "Mean growth ~0.3 mm·yr⁻¹ in sqrt space; range 0.1 → 0.5 by phenotype/genetics. [19]",
            "Function: BCVA poor; low-luminance VA & reading speed more sensitive.",
            "Perilesional FAF pattern (none/focal/banded/diffuse/patchy) prognoses growth (Holz FAM). [17]",
        ],
        eyebrow="03 · GA · measurement",
    )

    d.table_slide(
        "Approved drugs for GA (as of 2026)",
        ["Drug", "Target", "Trials", "Effect on growth", "Notes"],
        [
            ["Pegcetacoplan (Syfovre)", "C3", "OAKS, DERBY",
             "~14–22% reduction over 24 mo (sqrt scale)",
             "Apellis · monthly or bimonthly intravitreal · CNV conversion risk"],
            ["Avacincaptad pegol (Izervay)", "C5", "GATHER1, GATHER2",
             "~14–18% reduction over 12 mo",
             "Astellas/Iveric · monthly intravitreal · ocular inflammation rare"],
        ],
        col_widths=[2.7, 1.0, 1.7, 3.4, 3.5],
        eyebrow="03 · GA · therapeutics",
        footnote="Both inhibit the complement cascade. Both slow but do not stop atrophy. Neither shows functional vision improvement in pivotal trials. Refs [22], [23].",
    )

    d.callout(
        "Where the field is stuck",
        [
            "Structure-function disconnect: drugs slow lesion growth but do not improve vision in trials.",
            "Slow trials: 12–24 months × thousands of patients × monthly injections — costly and burdensome.",
            "Heterogeneous progression: same baseline area, very different 12-mo outcomes.",
            "Inadequate individualization: clinicians treat by lesion phenotype, not by patient-specific growth model.",
            "Limited synthetic-control precedent in ophthalmology vs oncology.",
        ],
        "All four problems map directly onto things a digital twin is good at — predicting heterogeneity, simulating interventions, and supplying virtual control arms.",
        eyebrow="03 · GA · unmet need",
    )

    # ====== SECTION 4 — WHY A GA TWIN ======
    d.section_divider(4, "Why a GA twin?",
                      "The clinical questions and the structural fit.")

    d.content(
        "Clinical questions a GA twin can answer",
        [
            "Personalized prognosis — at this baseline, what's the expected lesion area at 12, 24, 36 months, with uncertainty?",
            "Treatment scheduling — when is it worth starting therapy? Monthly or bimonthly? When can we stop?",
            "Risk-stratified enrichment — which patients should be enrolled in a fast-progressor trial cohort?",
            "Counterfactual reasoning — what would this patient look like at 24 mo without therapy?",
            "Trial design — virtual sample-size calculations under realistic patient heterogeneity and dropout.",
            "Synthetic control augmentation — supply twin-based controls to shrink the real control arm.",
            "Patient communication — show projected fundus/lesion + functional impact under treatment vs no treatment.",
        ],
        eyebrow="04 · Why",
    )

    d.callout(
        "Why GA fits a twin unusually well",
        [
            "Quantitative, validated structural endpoint (FAF lesion area).",
            "Mechanism partly known (complement-mediated RPE death) → enables hybrid mechanistic + ML twins, not pure black-box.",
            "Slow progression and 12–24 mo trial timelines make synthetic-control augmentation valuable.",
            "Strong genetics (CFH, ARMS2, C3) for personalization priors.",
            "Two FDA-approved drugs with modest, quantified effect sizes — the twin has clear intervention literature to calibrate against.",
            "Existing retinal foundation model (RETFound) provides imaging perception layer for free.",
        ],
        "Of all ophthalmology indications, GA is the strongest candidate to demonstrate a clinically and regulatorily useful patient digital twin.",
        eyebrow="04 · Why",
    )

    # ====== SECTION 5 — PRIOR ART ======
    d.section_divider(5, "Prior art in ophthalmology AI",
                      "Imaging progression models, retinal foundation models, in-silico work.")

    d.content(
        "Imaging-based GA progression prediction",
        [
            "Pfau M, Schmidt-Erfurth U et al. (Bonn / Vienna) — deep learning on FAF predicts 12-month GA growth, beats linear sqrt models; identifies perilesional features as drivers. [24]",
            "Niu S et al. — deep learning on OCT predicts AMD/GA conversion. [25]",
            "Schmidt-Erfurth U lab — automated AI quantification of OCT biomarkers (drusen, RPE atrophy, EZ integrity, hyperreflective foci) with prognostic value. [26]",
            "Liefers B et al. (Nijmegen / Moorfields) — deep segmentation of GA on FAF and OCT, replacing reading-center grading. [27]",
            "Yehoshua Z et al. — long-running natural-history series on FAF growth. [28]",
            "Common limitation: published work is pure ML prediction — no intervention simulator, no closed-loop twin.",
        ],
        eyebrow="05 · Prior art · imaging ML",
    )

    d.two_column(
        "Foundation models & twin precedents",
        "Retinal foundation models",
        [
            "RETFound (Zhou et al., Nature 2023) — open-weight ViT pretrained on ~1.6 M retinal images. State-of-the-art on AMD progression and several systemic-disease tasks. [1]",
            "Closed commercial counterparts: EyeFM, EyeBERT, RetinAI's pipeline.",
            "Self-supervised pretraining on UK Biobank retina shows representations carry information about systemic disease.",
        ],
        "In-silico / twin work touching ophthalmology",
        [
            "Unlearn.AI ophthalmology pilots — twin-based control augmentation in dry AMD (announced; trial-stage). [15]",
            "Mechanistic models of RPE death — reaction-diffusion equations for lesion expansion (Shen, Friedman). [29]",
            "Population-level disease-progression models for AMD using NLMEM (Roche / Genentech) for trial design. [30]",
        ],
        eyebrow="05 · Prior art",
    )

    d.callout(
        "Visible gaps the proposed twin fills",
        [
            "No published patient-level twin combining imaging-derived state with a mechanistic growth model AND an intervention simulator.",
            "No widely-deployed clinician-facing what-if tool for GA.",
            "Limited synthetic-control-arm precedent for GA specifically.",
            "Open question: does treatment effect generalize across phenotypes (focality, perilesional pattern, genetics)? A twin lets you simulate this.",
            "No open, sharable PoC the community can iterate on.",
        ],
        "There is room for a small, well-validated, OPEN PoC twin — and the gap is precisely the patient-level simulation + intervention layer.",
        eyebrow="05 · Prior art · gaps",
    )

    # ====== SECTION 6 — ARCHITECTURE ======
    d.section_divider(6, "Proposed architecture",
                      "Five layers, open foundation, hybrid mechanistic + ML core.")

    d.architecture_slide("Reference architecture for the GA twin",
                          eyebrow="06 · Architecture")

    d.content(
        "Mechanistic core",
        [
            "Growth law in sqrt-area space:  √A(t) = √A(0) + k_i · t  (Feuer 2013). [18]",
            "Per-patient growth coefficient k_i drawn from a log-normal prior parameterized by:",
            ("baseline area A(0)", 1),
            ("focality (uni- vs multifocal)", 1),
            ("perilesional FAF pattern (none / focal / banded / diffuse / patchy) — banded & diffuse are fastest [17]", 1),
            ("genetics (CFH, ARMS2 risk alleles), age, sex", 1),
            "Drug effect: multiplicative reduction on k_i (~0.18 pegcetacoplan, ~0.16 avacincaptad over 24 mo) [22][23], with patient heterogeneity.",
            "Optional reaction-diffusion extension for spatial lesion shape (Shen/Friedman). [29]",
        ],
        eyebrow="06 · Architecture · core",
    )

    d.content(
        "ML personalization layer + LLM orchestrator",
        [
            "ML personalization: gradient-boosted regressor (or small CNN on FAF/OCT) maps baseline features → posterior over k_i. Calibrated via isotonic / Platt.",
            "Optional: replace the imaging head with RETFound embeddings + a light task head — buys data efficiency for free. [1]",
            "Uncertainty: Bayesian hierarchical model OR ensemble + conformal prediction — produce 95% CIs on projected area at any horizon.",
            "LLM orchestrator (clinician surface): tool-using LLM calls into the twin and a knowledge base (FHIR + GA guideline corpus); never decides on its own. Aligned with FDA SaMD and EMA reflection paper. [12]",
            "Audit log: every clinician question, twin call, and answer is recorded for QA and for regulatory submission.",
        ],
        eyebrow="06 · Architecture · ML & LLM",
    )

    # ====== SECTION 7 — POC ======
    d.section_divider(7, "Proof of concept",
                      "What gets built first — scope, data, methodology, validation, use cases.")

    d.callout(
        "PoC v0 — scope",
        [
            "Patient-level GA twin with sqrt-area growth law and ML personalization — one disease, one endpoint, one drug class.",
            "Streamlit dashboard mirroring the cardiovascular twin already in this repo.",
            "In-silico mini-trial module — simulate N virtual patients, randomize, estimate power for detecting a given drug effect at 12/24 mo.",
            "Built on a synthetic cohort calibrated to PUBLISHED distributions (no credentialed data needed for v0).",
            "Roadmap escalates to AREDS2 (dbGaP), RETFound embeddings, LLM orchestrator.",
        ],
        "Tractable in weeks. Defensible without credentialed data. Extensible to a real-data v1.",
        eyebrow="07 · PoC · scope",
    )

    d.table_slide(
        "Data plan — five tiers, accessible to credentialed",
        ["Tier", "Source", "What you get", "Access"],
        [
            ["0  Synthetic", "Simulated cohort from published distributions",
             "Tabular longitudinal GA areas + covariates",
             "Immediate, no DUC — best for v0"],
            ["1  Public summary", "OAKS/DERBY, GATHER1/2, AREDS2 publications",
             "Calibration anchors (means, variance, drug effect)",
             "Free, manual extraction"],
            ["2  Open imaging", "STARE, EyePACS, Kaggle AMD; RETFound weights",
             "Cross-sectional retinal images + foundation embeddings",
             "Free"],
            ["3  Controlled access", "AREDS / AREDS2 via NEI dbGaP",
             "Real longitudinal AMD/GA cohort (~4500 pts, up to 12 yr)",
             "dbGaP DUC, ~weeks"],
            ["4  Clinical", "Health-system FAF/OCT under IRB",
             "Site-specific real twins, prospective forecasts",
             "IRB-dependent"],
        ],
        col_widths=[1.6, 3.2, 4.0, 3.5],
        eyebrow="07 · PoC · data",
        footnote="v0 sits at tiers 0+1. v1 escalates to tier 3. Tier 4 is pilot-deployment territory.",
    )

    d.content(
        "Required virtual-patient features",
        [
            "Demographics — age, sex.",
            "Baseline imaging-derived — GA area (mm² and √mm²), focality (uni/multi), perilesional FAF pattern (none/focal/banded/diffuse/patchy), reticular pseudodrusen present/absent.",
            "Optional genetics — CFH risk allele, ARMS2/HTRA1 risk allele.",
            "Functional — BCVA, low-luminance VA, optionally reading speed.",
            "Longitudinal observations — per-visit GA area at 0, 6, 12, 18, 24 months (configurable).",
            "Treatment record — drug, dose, schedule, adherence.",
            "Outcome derived during simulation — projected lesion area trajectory + uncertainty band; functional projection (optional).",
        ],
        eyebrow="07 · PoC · features",
    )

    d.content(
        "Methodology — eight steps for v0",
        [
            "1. Calibrate growth-rate distribution to literature: log-normal k_i with mean and variance by perilesional pattern, focality, baseline area.",
            "2. Sample a virtual cohort (e.g., 500 patients) with covariate distributions matching AREDS2 summaries.",
            "3. Simulate per-patient trajectories in sqrt space with measurement noise.",
            "4. Train a tabular ML model (gradient-boosted trees) on the synthetic cohort to recover k_i from baseline covariates; report calibration.",
            "5. Wrap in twin/ modules mirroring the cardiovascular twin (Patient, LesionGrowthModel, RiskModel, Simulator).",
            "6. Streamlit dashboard: patient picker, baseline state, projected lesion area at 12/24 mo with 95% bands, drug toggle, in-silico-trial tab.",
            "7. Sanity checks: growth-rate distribution recovery; simulated OAKS-like trial recovers the published effect at trial N.",
            "8. Document limitations clearly in UI and README; lock the v0 release.",
        ],
        eyebrow="07 · PoC · methodology",
    )

    d.two_column(
        "Validation plan",
        "v0 (synthetic only)",
        [
            "Distribution-recovery test: ML model recovers planted k_i distribution within tolerance.",
            "Trial-recovery test: simulated OAKS-like trial recovers published ~18% effect at trial N=600.",
            "Calibration plot for predicted area at 12 / 24 mo (pp-plot).",
            "Sensitivity analyses on prior parameters and noise levels.",
        ],
        "v1 (with credentialed data)",
        [
            "Internal: held-out AREDS2 patients; RMSE on 12/24 mo lesion area.",
            "External: a different cohort if available (Holekamp Proxima).",
            "Calibration: predicted vs observed area distribution.",
            "Drug-effect sanity: estimated effect should fall in OAKS/GATHER range.",
            "Bias audit: performance by baseline area decile, sex, age, ancestry.",
            "Prospective forecasting study at a clinical site (pilot).",
        ],
        eyebrow="07 · PoC · validation",
    )

    d.content(
        "Use cases the twin enables",
        [
            "Treatment scheduling — decide start, switch, intensify, hold, stop, given expected residual benefit.",
            "Trial design — virtual sample-size calculations under heterogeneity and dropout; enrichment for fast progressors.",
            "Synthetic control augmentation — supply twin-based controls for external-control-arm submissions (with regulator pre-alignment).",
            "Counterfactual reasoning — 'what would this patient look like at 24 mo without therapy?' as a decision aid.",
            "Patient communication — visualize projected fundus and functional impact under treatment vs no treatment.",
            "Drug-development hypothesis testing — what effect size on growth rate is needed for a meaningful functional endpoint?",
            "Education — train residents and trial investigators on lesion-progression intuition.",
        ],
        eyebrow="07 · PoC · use cases",
    )

    # ====== SECTION 8 — ROADMAP & RISKS ======
    d.section_divider(8, "Roadmap, risks, references",
                      "What ships when, what could break it, who said what.")

    d.table_slide(
        "Roadmap",
        ["Sprint", "Duration", "Deliverable"],
        [
            ["1", "1–2 weeks",  "v0 simulated GA twin + Streamlit dashboard (this PoC)"],
            ["2", "+1 month",   "RETFound embeddings on a public retinal image set; couple to growth-rate prior"],
            ["3", "+2 months",  "Apply for AREDS2 dbGaP, ingest, recalibrate, validate (v1)"],
            ["4", "+3 months",  "In-silico-trial module + LLM orchestrator with eval harness"],
            ["5", "+6 months",  "Pilot with a clinical site — IRB, prospective forecasting study"],
            ["6", "+9 months",  "Pre-submission meeting with FDA / EMA on synthetic-control-arm use"],
        ],
        col_widths=[1.0, 1.8, 9.5],
        eyebrow="08 · Roadmap",
    )

    d.callout(
        "Risks, ethics, regulatory",
        [
            "Synthetic data ≠ real data — v0 conclusions must never be used clinically; clearly labeled in UI and README.",
            "Bias — AREDS2 is European-ancestry-skewed; growth models may not generalize across ancestry; bias audit mandatory.",
            "Regulatory — clinical-decision use is SaMD; in-silico-trial use falls under MIDD; engage FDA/EMA early.",
            "Privacy — AREDS2/dbGaP carries strong DUC obligations; clinical data needs IRB and deidentification; OMOP/FHIR ingestion at the boundary.",
            "LLM-orchestrator risks — prompt injection, hallucinated tool calls, citation drift; eval harness and red-team before deployment.",
            "Generative components can leak training data and under-represent rare phenotypes; privacy testing + rare-class oversampling.",
        ],
        "Ship a clearly-labeled teaching PoC first. Build credibility with validation. Engage regulators before any clinical claim.",
        eyebrow="08 · Risks",
    )

    # ====== References ======
    d.content(
        "References — modern DT & methodology",
        [
            "[1] Zhou Y et al. A foundation model for generalizable disease detection from retinal images. Nature 2023.",
            "[2] Singhal K et al. Towards Expert-Level Medical Question Answering with Med-PaLM 2. arXiv 2305.09617 / Nature 2023.",
            "[3] Moor M et al. Foundation models for generalist medical AI. Nature 2023.",
            "[4] Google. MedGemma technical report. 2024.",
            "[5] Tu T et al. Towards Conversational Diagnostic AI (AMIE). DeepMind, 2024.",
            "[6] Chen RTQ et al. Neural Ordinary Differential Equations. NeurIPS 2018.",
            "[7] Raissi M, Perdikaris P, Karniadakis GE. Physics-informed neural networks. JCP 2019.",
            "[8] Rackauckas C et al. Universal Differential Equations. 2020.",
            "[9] Kazerouni A et al. Diffusion models in medical imaging — survey. MIA 2023.",
            "[10] Xu L et al. CTGAN. NeurIPS 2019; Borisov V et al. Tabular generative models survey. 2023.",
            "[11] Theodorou B et al. Synthesize EHR with generative AI. npj Digital Medicine 2023.",
            "[12] FDA AI/ML SaMD Action Plan 2021; EMA reflection paper on AI in medicines lifecycle 2024.",
            "[13] FDA Model-Informed Drug Development (MIDD) Pilot Program.",
            "[14] Thorlund K et al. Synthetic and external controls. Clin Epidemiol 2020.",
            "[15] Unlearn.AI — TwinRCT / Procova.",
        ],
        eyebrow="08 · References (1/2)",
    )

    d.content(
        "References — GA & ophthalmology",
        [
            "[16] Wong WL et al. Global prevalence of AMD. Lancet Glob Health 2014.",
            "[17] Holz FG et al. FAM Study — perilesional FAF patterns and GA progression. Am J Ophthalmol 2007.",
            "[18] Feuer WJ et al. Square-root transformation of GA area. Ophthalmology 2013.",
            "[19] Yehoshua Z et al. Natural history of GA in AMD. Ophthalmology 2011.",
            "[20] Chew EY et al. AREDS2: A randomized clinical trial. JAMA 2013.",
            "[21] Holekamp NM et al. PROXIMA A and B natural history of GA. Ophthalmol Retina 2020.",
            "[22] Heier JS et al. Pegcetacoplan for GA — OAKS and DERBY. Lancet 2023.",
            "[23] Khanani AM et al. Avacincaptad pegol for GA — GATHER2. Lancet 2023.",
            "[24] Pfau M et al. Deep learning predicts GA progression on FAF. IOVS / Ophthalmology 2020–2022.",
            "[25] Niu S et al. Deep learning prediction of progression to AMD from OCT. Ophthalmol Sci 2020.",
            "[26] Schmidt-Erfurth U et al. AI in retina. Prog Retin Eye Res 2018–2024.",
            "[27] Liefers B et al. Deep segmentation of GA on FAF/OCT. Ophthalmology / TVST 2020.",
            "[28] Yehoshua Z et al. FAF imaging in GA — natural history series. AJO / Retina.",
            "[29] Shen JK et al. Mathematical modeling of RPE death and lesion expansion.",
            "[30] Roche / Genentech population PK/PD reports for AMD trial design.",
        ],
        eyebrow="08 · References (2/2)",
    )

    # ====== Closing ======
    d.callout(
        "Closing & questions",
        [
            "Modern medical twins = foundation-model perception + hybrid mechanistic/ML simulation + LLM orchestration + closed clinical loop.",
            "GA is the right ophthalmology beachhead — quantitative endpoint, slow trials, two new drugs with measurable effects, mechanistic intuition.",
            "v0 PoC is buildable on synthetic data calibrated to literature; v1 is a 2-month escalation to AREDS2.",
            "Roadmap: PoC → AREDS2 → in-silico trials → clinical pilot → regulator pre-submission.",
            "What's the strongest objection? What clinical partner do we want to engage first?",
        ],
        "Open questions for review: choice of clinical partner, dbGaP timing, scope of v0 vs v1, regulator engagement plan.",
        eyebrow="Discussion",
    )

    OUT.parent.mkdir(exist_ok=True)
    d.prs.save(OUT)
    print(f"[deck] wrote {OUT}  ({OUT.stat().st_size/1024:.1f} KB, "
          f"{len(d.prs.slides)} slides)")


if __name__ == "__main__":
    build()
